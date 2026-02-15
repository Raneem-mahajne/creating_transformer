"""
Create a PowerPoint presentation for the plus_last_even configuration.
Run from project root: python scripts/create_plus_last_even_presentation.py
"""
import os
import sys

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt


def load_config():
    config_path = os.path.join(
        os.path.dirname(os.path.dirname(__file__)),
        "configs", "plus_last_even.yaml"
    )
    with open(config_path, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, item in enumerate(bullets):
        p = body.add_paragraph() if i == 0 else body.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
    return slide


def add_two_column_slide(prs, title, left_bullets, right_bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    for i, item in enumerate(left_bullets):
        p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(4)

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    for i, item in enumerate(right_bullets):
        p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(4)
    return slide


def add_code_slide(prs, title, code_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(24)
    title_box.text_frame.paragraphs[0].font.bold = True

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
    tf = code_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.space_after = Pt(2)
    return slide


def main():
    config = load_config()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- Research presentation: How the transformer implements the rule ----

    # Slide 1: Title
    add_title_slide(
        prs,
        "How a Transformer Implements the Plus-Last-Even Rule",
        "Interpretable minimal transformers: geometry as algorithm"
    )

    # Slide 2: The task
    add_content_slide(prs, "The Plus-Last-Even Task", [
        "Rule: When the current token is '+', the next token must be the most recent even number that appeared before that '+'.",
        "",
        "Example:  5, 3, 8, 7, +, 8,   10, 2, 4, +, 4, ...",
        "                         ↑ last even = 8         ↑ last even = 4",
        "",
        "The model must: (1) detect that the current position is '+', (2) search backward to find the last even number, (3) output that number.",
        "This is operator-conditional and context-dependent — a good testbed for understanding how transformers perform retrieval."
    ])

    # Slide 3: Transformer pipeline (high level)
    add_content_slide(prs, "Decoder-Only Transformer: High-Level Pipeline", [
        "Input: sequence of token IDs → each position gets a hidden state in ℝ^d.",
        "",
        "1. Embedding: Token embedding E_t + position embedding P_i → x_i = E_{token(i)} + P_i  (each (token, position) is a point in ℝ^d).",
        "2. Self-attention: For each position i, compute query q_i = W_Q x_i, key k_i = W_K x_i, value v_i = W_V x_i; attention weights α_ij ∝ exp(q_i·k_j/√d) (causal: j≤i); output at i = Σ_j α_ij v_j.",
        "3. Residual: state_i := x_i + attention_output_i.",
        "4. Feedforward: state_i := state_i + FFN(state_i).",
        "5. LM head: logits = state W_lm^T + b → softmax → P(next token).",
        "",
        "Next-token prediction at position i uses only positions 0..i; the model learns to use attention to 'look back' when needed."
    ])

    # Slide 4: Why 2D? Interpretability
    add_content_slide(prs, "Why We Use 2D Embeddings (n_embd=2, d_k=2)", [
        "In full-scale transformers, d is large (e.g. 768, 4096); internal states are opaque.",
        "We set n_embd = 2 and head size = 2 so that every vector lives in ℝ²:",
        "  • Embeddings, queries, keys, values, attention outputs, and pre-LM-head states are all 2D points.",
        "  • No PCA or t-SNE — what we plot is the actual model state.",
        "Claim: The learned 2D geometry is the algorithm. The arrangement of points and decision boundaries in the plane can be read as the step-by-step procedure the model executes.",
        "",
        "Config: block_size=8, 1 head; single block (embed → attention → residual → FFN → LM head)."
    ])

    # Slide 5: How the transformer implements the rule — four steps
    add_content_slide(prs, "How the Transformer Implements Plus-Last-Even (4 Steps)", [
        "(1) Representation: Token and position embeddings assign each (token, position) a unique point in ℝ². So '8 at position 3' and '+' at position 5 have distinct locations; the model 'sees' who is where.",
        "",
        "(2) Retrieval: W_Q and W_K map embeddings to queries and keys. The model learns so that the query at a '+' position is close (high dot product) to the keys of even-number positions, with recency encoded — the key of the most recent even gets the highest attention.",
        "",
        "(3) Read-out: Attention weights α_ij select the relevant positions; the value vectors at those positions are summed. So at '+', the attention output is dominated by the value at the 'last even' position. The residual then adds this to the current embedding: state moves in the plane toward the 'last even' representation.",
        "",
        "(4) Output: The LM head is a linear map + softmax; it partitions ℝ² into regions (one per output token). The final state (after residual + FFN) lands in one region → that token is predicted. For correct behavior, the state at '+' must land in the region of the last even number."
    ])

    # Slide 6: Step 1 — Representation
    add_content_slide(prs, "Step 1: Representation (Token + Position in ℝ²)", [
        "x_i = E_{token(i)} + P_i. Each (token, position) pair is a unique point in the plane.",
        "Visualization: embeddings.png, token_position_embedding_space.png — scatter of all (token, position) points.",
        "The model can distinguish '6 at position 2' from '6 at position 5' and from '+' at position 4 because they lie at different 2D locations."
    ])

    # Slide 7: Step 2 — Retrieval (Q/K)
    add_content_slide(prs, "Step 2: Retrieval — Queries and Keys", [
        "Attention weight of position i on j: α_ij ∝ exp(q_i · k_j / √d). So 'who attends to whom' is determined by query–key alignment.",
        "For plus-last-even, the query at a '+' position must align with keys of even-number positions; among those, the key of the most recent even should have the largest dot product with the '+' query.",
        "Visualization: qk_embedding_space.png (scatter of all Q and K in ℝ²), qk_full_attention_heatmap.png (attention matrix over token–position pairs). We can see that '+' queries attend strongly to the 'last even' key."
    ])

    # Slide 8: Step 3 — Read-out (Values + Residual)
    add_content_slide(prs, "Step 3: Read-out — Values and Residual", [
        "Attention output at position i: Σ_j α_ij v_j. So the state receives a weighted sum of value vectors from previous positions.",
        "At a '+' position, α is concentrated on the 'last even' position → the sum is essentially the value at that position. The residual then does: new_state = x_i + (attention output). So the state moves in the plane from '+' toward the 'last even' representation.",
        "Visualization: residuals.png (embed vs V_transformed vs final per position), v_before_after_demo_*.png (trajectory of each position's state overlaid on P(next) heatmaps; green = correct prediction, red = wrong)."
    ])

    # Slide 9: Step 4 — Output (LM Head)
    add_content_slide(prs, "Step 4: Output — LM Head Partitions ℝ²", [
        "After residual and feedforward, the state is a point in ℝ². The LM head applies a linear map + softmax to produce P(next token).",
        "So the plane is partitioned into regions: if the state falls in region R_k, the model predicts token k. Decision boundaries are linear.",
        "Visualization: lm_head_probability_heatmaps.png, probability_heatmap_with_embeddings.png, probability_heatmap_with_values.png — heatmaps of P(next = token) over the 2D space; overlays show where each position's state lands relative to these regions."
    ])

    # Slide 10: Experimental setup (brief)
    data = config["data"]
    model = config["model"]
    training = config["training"]
    add_content_slide(prs, "Experimental Setup (plus_last_even config)", [
        "Data: PlusLastEvenRule; integers 0–10 and '+'; 2000 sequences, length 20–50; operator probability 0.3.",
        "Model: n_embd=2, block_size=8, 1 head, head_size=2.",
        "Training: 20k steps, batch 8, lr=0.001; checkpoints every 100 (for videos and step-wise visualization).",
        "All figures: plus_last_even/plots/. Run: python main.py plus_last_even [--visualize | --video]."
    ])

    # Slide 11: Summary / takeaway
    add_content_slide(prs, "Summary", [
        "Plus-last-even: after '+', output the most recent even number. The transformer implements this by: (1) representing (token, position) in ℝ², (2) aligning '+' queries with 'last even' keys so attention retrieves the right position, (3) adding that value via the residual so the state moves to the 'last even' region, (4) mapping that state to the correct token via the LM head.",
        "By constraining to 2D, the learned geometry is the algorithm — we can read off each step from the visualizations.",
        "This framework generalizes to other procedural rules (copy-modulo, successor, plus-max-of-two, etc.) with the same interpretability pipeline."
    ])

    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "plus_last_even")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "plus_last_even_presentation.pptx")
    prs.save(out_path)
    print(f"Presentation saved to {out_path}")


if __name__ == "__main__":
    main()
